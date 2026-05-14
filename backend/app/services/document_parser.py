from pathlib import Path
from io import BytesIO
import csv

import fitz
import pdfplumber
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image


class ParsedDocument:
    def __init__(self, text: str, hidden_segments: list[str] | None = None):
        self.text = text
        self.hidden_segments = hidden_segments or []


def parse_document_bytes(filename: str, content: bytes) -> ParsedDocument:
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(content)
    if ext == ".docx":
        return _parse_docx(content)
    if ext == ".pptx":
        return _parse_pptx(content)
    if ext == ".xlsx":
        return _parse_xlsx(content)
    if ext == ".csv":
        return _parse_csv(content)
    if ext in {".txt", ".md"}:
        return ParsedDocument(content.decode("utf-8", errors="ignore"))
    if ext in {".html", ".htm"}:
        return _parse_html(content)
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".tiff", ".tif"}:
        return ParsedDocument("")

    return ParsedDocument(content.decode("utf-8", errors="ignore"))


def _parse_pdf(content: bytes) -> ParsedDocument:
    visible: list[str] = []
    hidden: list[str] = []

    pdf = fitz.open(stream=content, filetype="pdf")
    for page in pdf:
        blocks = page.get_text("dict").get("blocks", [])
        for block in blocks:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    txt = span.get("text", "").strip()
                    if not txt:
                        continue
                    color = span.get("color", 0)
                    size = float(span.get("size", 12))
                    if size <= 1.5 or color == 0xFFFFFF:
                        hidden.append(txt)
                    else:
                        visible.append(txt)

    with pdfplumber.open(BytesIO(content)) as plumber_pdf:
        for page in plumber_pdf.pages:
            extracted = page.extract_text() or ""
            if extracted.strip():
                visible.append(extracted)

    return ParsedDocument("\n".join(visible), hidden)


def _parse_docx(content: bytes) -> ParsedDocument:
    doc = DocxDocument(BytesIO(content))
    lines = [p.text for p in doc.paragraphs if p.text.strip()]
    return ParsedDocument("\n".join(lines))


def _parse_pptx(content: bytes) -> ParsedDocument:
    prs = Presentation(BytesIO(content))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return ParsedDocument("\n".join(texts))


def _parse_xlsx(content: bytes) -> ParsedDocument:
    wb = load_workbook(filename=BytesIO(content), data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = " | ".join(str(cell) for cell in row if cell is not None)
            if line:
                rows.append(line)
    return ParsedDocument("\n".join(rows))


def _parse_csv(content: bytes) -> ParsedDocument:
    decoded = content.decode("utf-8", errors="ignore")
    reader = csv.reader(decoded.splitlines())
    lines = [" | ".join(row) for row in reader if row]
    return ParsedDocument("\n".join(lines))


def _parse_html(content: bytes) -> ParsedDocument:
    soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")

    hidden_parts: list[str] = []
    for tag in soup.find_all(style=True):
        style = tag.attrs.get("style", "").lower()
        if "display:none" in style or "visibility:hidden" in style or "font-size:0" in style:
            hidden_text = tag.get_text(" ", strip=True)
            if hidden_text:
                hidden_parts.append(hidden_text)

    for script in soup(["script", "style"]):
        script.decompose()

    text = soup.get_text("\n", strip=True)
    return ParsedDocument(text, hidden_parts)


def load_pillow_image(content: bytes) -> Image.Image:
    return Image.open(BytesIO(content))

